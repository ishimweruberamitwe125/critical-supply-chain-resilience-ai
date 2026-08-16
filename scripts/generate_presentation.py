"""Generate a PowerPoint presentation for researchers, companies, and grant funders."""

from __future__ import annotations

import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

IMAGES_DIR = PROJECT_ROOT / "docs" / "images"
OUTPUT_DIR = PROJECT_ROOT / "docs" / "presentations"
OUTPUT_FILE = OUTPUT_DIR / "Critical_Supply_Chain_Resilience_AI.pptx"

ACCENT = RGBColor(37, 99, 235)
DARK = RGBColor(15, 23, 42)
GRAY = RGBColor(71, 85, 105)


def _set_title(text_frame, text: str, size: int = 36) -> None:
    text_frame.text = text
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(size)
    paragraph.font.bold = True
    paragraph.font.color.rgb = DARK


def _add_bullets(text_frame, bullets: list[str], size: int = 20) -> None:
    text_frame.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = GRAY


def _add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.5))
    tf = box.text_frame
    tf.text = "Critical Supply Chain Resilience AI"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.5))
    sub_tf = sub.text_frame
    sub_tf.text = (
        "AI-Powered Decision Support for Critical Manufacturing Supply Chains\n"
        "Prototype Overview for Researchers, Industry, and Funders\n"
        "David Ishimwe Ruberamitwe"
    )
    for paragraph in sub_tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = GRAY


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.8))
    _set_title(title_box.text_frame, title, size=30)

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.5))
    _add_bullets(body.text_frame, bullets)


def _add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.7))
    _set_title(title_box.text_frame, title, size=28)

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.2), width=Inches(11.8))

    if caption:
        cap = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5))
        cap.text_frame.text = caption
        cap.text_frame.paragraphs[0].font.size = Pt(14)
        cap.text_frame.paragraphs[0].font.color.rgb = GRAY


def build_presentation() -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs)

    _add_bullet_slide(
        prs,
        "Who This Project Is For",
        [
            "Researchers — reproducible supply chain resilience analytics and benchmark networks.",
            "Companies — supplier risk visibility, disruption stress testing, and mitigation planning.",
            "Grant funders & policymakers — evidence-based view of critical infrastructure vulnerability.",
            "Open platform designed for collaboration, pilots, and sector-specific extensions.",
        ],
    )

    _add_bullet_slide(
        prs,
        "Business Value for Companies",
        [
            "Identify single-source suppliers before a disruption becomes a production crisis.",
            "Quantify service-level impact of factory outages, supplier failures, and logistics delays.",
            "Compare mitigation options using simulation and ranked diversification recommendations.",
            "Reduce uncertainty in strategic sourcing, inventory planning, and continuity decisions.",
        ],
    )

    _add_bullet_slide(
        prs,
        "The Challenge",
        [
            "Critical manufacturing supply chains are globally interconnected and fragile.",
            "Semiconductor, energy, medical, aerospace, and defense disruptions have national impact.",
            "Organizations lack quantitative tools to anticipate shocks and compare resilience strategies.",
        ],
    )

    _add_bullet_slide(
        prs,
        "Research Gap",
        [
            "Limited open frameworks combining graph analytics, ML prediction, simulation, and optimization.",
            "Existing tools often focus on cost efficiency rather than resilience under disruption.",
            "Need for reproducible research platforms supporting policy and strategic decision-making.",
        ],
    )

    _add_bullet_slide(
        prs,
        "Project Vision",
        [
            "Build an open AI research and engineering platform for supply chain resilience.",
            "Detect early disruption signals and quantify vulnerability across supplier networks.",
            "Simulate what-if scenarios and recommend diversification and logistics strategies.",
            "Support researchers, manufacturers, and policymakers with evidence-based analytics.",
        ],
    )

    _add_bullet_slide(
        prs,
        "System Architecture",
        [
            "Data ingestion: supplier, facility, and material-flow datasets.",
            "Graph model: NetworkX representation of suppliers, factories, and distribution hubs.",
            "Analytics pipeline: metrics → prediction → simulation → optimization → reporting.",
            "Interfaces: CLI demo, Jupyter case studies, and interactive Streamlit dashboard.",
        ],
    )

    _add_bullet_slide(
        prs,
        "Core Capabilities",
        [
            "Disruption Prediction — ML models on supplier reliability and geopolitical risk signals.",
            "Resilience Metrics — risk index, dependency scores, propagation, recovery time.",
            "Supply Chain Simulation — factory shutdowns, supplier outages, trade restrictions.",
            "Network Optimization — mitigation recommendations and logistics route planning.",
        ],
    )

    _add_bullet_slide(
        prs,
        "Methodology",
        [
            "Synthetic benchmark networks for semiconductor and energy infrastructure sectors.",
            "Graph-based vulnerability analysis and scenario stress testing.",
            "Prototype ML classifier for supplier disruption risk tiers.",
            "Modular Python architecture designed for extension to real operational datasets.",
        ],
    )

    _add_bullet_slide(
        prs,
        "Prototype Scope (Current vs Planned)",
        [
            "Current: synthetic semiconductor and energy networks with full analytics pipeline.",
            "Current: CLI demo, Streamlit dashboard, Jupyter case studies, and resilience metrics.",
            "Planned: integration with company ERP, logistics, and macroeconomic data feeds.",
            "Planned: PyTorch prediction models and OR-Tools mathematical optimization.",
        ],
    )

    image_slides = [
        (
            "Interactive Prototype Dashboard",
            IMAGES_DIR / "web_dashboard_preview.png",
            "Streamlit dashboard with industry selector, interactive Plotly charts, and what-if simulation.",
        ),
        (
            "Executive Resilience Overview",
            IMAGES_DIR / "dashboard_overview.png",
            "Composite KPI view: risk index, supplier dependency, scenario impact, and mitigation priorities.",
        ),
        (
            "Supplier Disruption Risk",
            IMAGES_DIR / "supplier_risk_chart.png",
            "Supplier risk tiers help prioritize monitoring and alternate sourcing investments.",
        ),
        (
            "Simulation & Mitigation Insights",
            IMAGES_DIR / "simulation_impact.png",
            "Baseline vs disrupted throughput across supplier and factory failure scenarios.",
        ),
        (
            "Supply Network Topology",
            IMAGES_DIR / "supply_network_graph.png",
            "Graph view of suppliers, factories, and distribution hubs in the benchmark network.",
        ),
    ]

    for title, image_path, caption in image_slides:
        if image_path.exists():
            _add_image_slide(prs, title, image_path, caption)
        else:
            _add_bullet_slide(
                prs,
                title,
                [
                    caption,
                    f"Generate image assets with: python scripts/generate_dashboard_images.py",
                ],
            )

    _add_bullet_slide(
        prs,
        "Semiconductor Case Study",
        [
            "8 suppliers, 3 fabrication/assembly plants, 1 distribution hub.",
            "Identifies single-source dependencies for wafers, packaging, and rare earth materials.",
            "Simulates wafer and packaging outages with service-level and recovery estimates.",
        ],
    )

    _add_bullet_slide(
        prs,
        "Energy Infrastructure Case Study",
        [
            "Grid transformer and inverter supply network with 8 specialized suppliers.",
            "Highlights risks from permanent magnets, transformer steel, and control software.",
            "Supports diversification planning for critical energy deployment equipment.",
        ],
    )

    _add_bullet_slide(
        prs,
        "Policy & Strategic Impact",
        [
            "Strengthens national security through resilient critical supply networks.",
            "Improves economic stability during global shocks and shortages.",
            "Supports public health preparedness and technological competitiveness.",
            "Provides transparent analytics for research and policy evaluation.",
        ],
    )

    _add_bullet_slide(
        prs,
        "Future Research Directions",
        [
            "Supply chain digital twins with live logistics data integration.",
            "Reinforcement learning for adaptive logistics optimization.",
            "Multi-agent coordination across suppliers and manufacturers.",
            "Expanded sector models: medical devices, aerospace, and defense manufacturing.",
        ],
    )

    _add_bullet_slide(
        prs,
        "How to Run the Prototype",
        [
            "Install: pip install -r requirements.txt",
            "Prepare data: python scripts/prepare_week1_data.py && python scripts/prepare_energy_data.py",
            "Terminal demo: python examples/semiconductor_supply_chain_demo.py",
            "Web dashboard: streamlit run dashboard/app.py",
            "GitHub: github.com/ishimweruberamitwe125/critical-supply-chain-resilience-ai",
        ],
    )

    _add_bullet_slide(
        prs,
        "Collaboration & Next Steps",
        [
            "Industry pilots with anonymized supplier and logistics datasets.",
            "University research partnerships on resilience metrics and ML models.",
            "Grant-funded expansion to medical, aerospace, and defense supply networks.",
            "Contact: David Ishimwe Ruberamitwe — ishimwerubera@gmail.com",
        ],
    )

    closing = prs.slides.add_slide(prs.slide_layouts[6])
    box = closing.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10), Inches(2))
    tf = box.text_frame
    tf.text = "Thank You"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    sub = closing.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10), Inches(1))
    sub_tf = sub.text_frame
    sub_tf.text = "Critical Supply Chain Resilience AI\nQuestions & Discussion"
    for paragraph in sub_tf.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = GRAY
        paragraph.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT_FILE)
    return OUTPUT_FILE


def main() -> None:
    output = build_presentation()
    print(f"Presentation saved to: {output}")


if __name__ == "__main__":
    main()
